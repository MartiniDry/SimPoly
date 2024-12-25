from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


# Fonction pour extraire les points des polygones d'un fichier PowerPoint
def extract_polygon_points(pptx_file):
    prs = Presentation(pptx_file)
    polygons = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.FREEFORM:
                pts = []
                for sp in shape.element.spPr.xpath('.//a:pt'):
                    if sp.getparent().tag.endswith('lnTo') or sp.getparent().tag.endswith('moveTo'):
                        x = round(float(sp.get('x')) / 9525, 2)
                        y = round(float(sp.get('y')) / 9525, 2)
                        pts.append((x, y))

                if pts[0] == pts[-1]:  # Si le polygone est fermé en son bout, alors retirer le dernier point.
                    pts.pop()

                polygons.append(pts)

    return polygons


# Fonction pour sauvegarder les points des polygones dans un fichier
def save_polygon_points(polygons, output_file):
    with open(output_file, 'w') as f:
        for polygon in polygons:
            for point in polygon:
                f.write(f"{point[0]} {point[1]}\n")
            f.write("\n")


# Extraire les points des polygones du fichier PowerPoint
pptx_file = "C://Users//T0267715//Downloads//dessin_chat.pptm"
polygons = extract_polygon_points(pptx_file)

# Sauvegarder les points des polygones dans 'chat.spl'
output_file = "polygon_files//chat.spl"
save_polygon_points(polygons, output_file)

print(f"Les points des polygones ont été sauvegardés dans {output_file}.")